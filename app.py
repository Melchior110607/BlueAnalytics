import streamlit as st
import plotly.graph_objects as go
import os
from company_analyzer import CompanyAnalyzer
from openai import OpenAI

# Configuration de la page - DOIT être en premier
st.set_page_config(
    page_title="BlueAnalytics - Analyse d'Entreprises",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Chargement des variables d'environnement


def apply_bluesign_colors():
    """Application des couleurs Bluesign"""
    st.markdown("""
    <style>
    .main-header {
        background: linear-gradient(90deg, #1E3A8A 0%, #3B82F6 100%);
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        color: white;
        text-align: center;
    }
    
    .header-container {
        background: linear-gradient(90deg, #1E3A8A 0%, #3B82F6 100%);
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        color: white;
        text-align: center;
    }
    
    .main-title {
        font-size: 3rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
        color: white;
    }
    
    .subtitle {
        font-size: 1.2rem;
        opacity: 0.9;
        color: white;
    }
    
    .metric-card {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        border-left: 4px solid #3B82F6;
        margin-bottom: 1rem;
    }
    
    .metric-card-company {
        background: linear-gradient(135deg, #1E3A8A 0%, #3B82F6 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        border-left: 4px solid #60A5FA;
        margin-bottom: 1rem;
    }
    
    .metric-card-company h2 {
        color: white !important;
    }
    
    .metric-card-company h3 {
        color: white !important;
        opacity: 0.9;
    }
    
    .sustainability-metric {
        background: linear-gradient(135deg, #10B981 0%, #059669 100%);
        color: white;
        padding: 1rem;
        border-radius: 8px;
        margin: 0.5rem 0;
    }
    
    .stButton > button {
        background: linear-gradient(90deg, #1E3A8A 0%, #3B82F6 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.5rem 1rem;
        font-weight: bold;
    }
    
    .stButton > button:hover {
        background: linear-gradient(90deg, #1E40AF 0%, #2563EB 100%);
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }
    </style>
    """, unsafe_allow_html=True)

# Couleurs Bluesign
BLUESIGN_COLORS = {
    'primary': '#1E3A8A',
    'secondary': '#3B82F6', 
    'accent': '#60A5FA',
    'light': '#DBEAFE',
    'dark': '#1E40AF',
    'success': '#10B981',
    'warning': '#F59E0B',
    'error': '#EF4444'
}

# Configuration CSS personnalisée
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(90deg, #1E3A8A 0%, #3B82F6 100%);
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        color: white;
        text-align: center;
    }
    .metric-card {
        background: linear-gradient(135deg, #1E3A8A 0%, #3B82F6 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        border-left: 4px solid #60A5FA;
        margin: 1rem 0;
    }
    
    .metric-card-company {
        background: linear-gradient(135deg, #1E3A8A 0%, #3B82F6 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        border-left: 4px solid #60A5FA;
        margin: 1rem 0;
    }
    
    .metric-card-company h2 {
        color: white !important;
    }
    
    .metric-card-company h3 {
        color: white !important;
        opacity: 0.9;
    }
    .chart-container {
        background: white;
        padding: 1.5rem;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        margin: 1rem 0;
    }
    .stButton > button {
        background: linear-gradient(90deg, #1E3A8A 0%, #3B82F6 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.75rem 1.5rem;
        font-weight: 600;
        transition: all 0.3s ease;
    }
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 15px rgba(59, 130, 246, 0.3);
    }
    .sidebar .sidebar-content {
        background: linear-gradient(180deg, #1E3A8A 0%, #3B82F6 100%);
    }
    .sustainability-metric {
        background: linear-gradient(135deg, #10B981 0%, #059669 100%);
        color: white;
        padding: 1rem;
        border-radius: 8px;
        margin: 0.5rem 0;
    }
</style>
""", unsafe_allow_html=True)

def main():
    # Application des couleurs Bluesign
    apply_bluesign_colors()
    
    # Header principal
    st.markdown("""
    <div class="header-container">
        <h1 class="main-title">🏢 BlueAnalytics</h1>
        <p class="subtitle">Analyse complète d'entreprises avec métriques de durabilité</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Sidebar pour la navigation
    with st.sidebar:
        st.markdown("### 📈 Navigation")
        page = st.selectbox(
            "Choisir une section",
            ["🏢 Analyse d'Entreprise", "📄 Analyse de Rapport"]
        )
    
    # Navigation des pages
    if page == "🏢 Analyse d'Entreprise":
        show_company_analysis()
    elif page == "📄 Analyse de Rapport":
        show_report_analysis()

def show_company_analysis():
    """Page d'analyse d'entreprise"""
    
    st.markdown("## 🏢 Analyse d'Entreprise")
    
    # Formulaire de saisie
    with st.container():
        col1, col2 = st.columns([2, 1])
        
        with col1:
            company_name = st.text_input(
                "Nom de l'entreprise",
                placeholder="Ex: Apple, Tesla, Microsoft, Google...",
                help="Entrez le nom de l'entreprise à analyser"
            )
        
        with col2:
            st.markdown("<br>", unsafe_allow_html=True)
            analyze_button = st.button("🔍 Analyser", use_container_width=True)
    
    if analyze_button and company_name:
        with st.spinner("🔍 Analyse en cours avec recherche web..."):
            try:
                # Initialiser l'analyseur d'entreprise
                analyzer = CompanyAnalyzer()
                
                # Analyser l'entreprise
                company_data = analyzer.analyze_company(company_name)
                
                if company_data and not company_data.get('error'):
                    display_company_results(company_data, company_name)
                else:
                    error_msg = company_data.get('error', 'Impossible d\'analyser cette entreprise. Veuillez vérifier le nom.')
                    st.error(f"❌ {error_msg}")
                    
            except Exception as e:
                st.error(f"❌ Erreur lors de l'analyse: {str(e)}")
    
    elif analyze_button and not company_name:
        st.warning("⚠️ Veuillez entrer le nom d'une entreprise")

def show_report_analysis():
    """Page d'analyse de rapport multi-formats via GPT-5-mini"""
    st.markdown("## 📄 Analyse de Rapport (PDF/Doc/Image/Texte)")

    uploaded_file = st.file_uploader(
        "Téléversez un rapport (PDF, DOCX, PPTX, TXT, PNG, JPG)",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"],
        accept_multiple_files=False
    )

    instruction = st.text_area(
        "Instruction d'analyse",
        placeholder="Ex: Résume les points clés ESG, extrais les KPIs financiers, et propose 3 graphiques pertinents",
        height=120
    )

    col1, col2 = st.columns([1,1])
    with col1:
        run_btn = st.button("🧠 Lancer l'analyse", use_container_width=True)
    with col2:
        clear_btn = st.button("🧹 Réinitialiser", use_container_width=True)

    if clear_btn:
        st.experimental_rerun()

    if run_btn:
        if not uploaded_file and not instruction:
            st.warning("Veuillez fournir un fichier ou une instruction d'analyse.")
            return

        with st.spinner("Analyse du rapport en cours..."):
            try:
                client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

                # Préparation du contenu
                file_bytes = uploaded_file.read() if uploaded_file else None
                file_name = uploaded_file.name if uploaded_file else "instruction.txt"

                # Construire les blocs de contenu pour Responses API
                content_parts = []
                if instruction:
                    content_parts.append({"type": "input_text", "text": instruction})

                extracted_text = None
                if file_bytes:
                    import base64
                    file_ext = os.path.splitext(file_name)[1].lower()
                    if file_ext in [".png", ".jpg", ".jpeg"]:
                        # Envoyer l'image encodée
                        b64 = base64.b64encode(file_bytes).decode("utf-8")
                        mime = f"image/{'jpeg' if file_ext in ['.jpg', '.jpeg'] else 'png'}"
                        content_parts.append({
                            "type": "input_image",
                            "image_data": b64,
                            "mime_type": mime
                        })
                        content_parts.append({"type": "input_text", "text": f"Image transmise: {file_name}. Analyse et extrais les informations pertinentes."})
                    elif file_ext == ".txt":
                        try:
                            extracted_text = file_bytes.decode("utf-8", errors="ignore")
                        except Exception:
                            extracted_text = None
                    elif file_ext == ".pdf":
                        # Tentative d'extraction de texte PDF si lib dispo
                        try:
                            import io
                            try:
                                from PyPDF2 import PdfReader
                                reader = PdfReader(io.BytesIO(file_bytes))
                                pages_text = []
                                for p in reader.pages:
                                    pages_text.append(p.extract_text() or "")
                                extracted_text = "\n\n".join(pages_text)
                            except Exception:
                                extracted_text = None
                        except ImportError:
                            extracted_text = None
                    elif file_ext == ".docx":
                        try:
                            from docx import Document
                            import io
                            doc = Document(io.BytesIO(file_bytes))
                            paragraphs = [para.text for para in doc.paragraphs]
                            extracted_text = "\n".join(paragraphs)
                        except Exception:
                            extracted_text = None
                    elif file_ext == ".pptx":
                        try:
                            from pptx import Presentation
                            import io
                            prs = Presentation(io.BytesIO(file_bytes))
                            texts = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        texts.append(shape.text)
                            extracted_text = "\n".join(texts)
                        except Exception:
                            extracted_text = None

                # Ajouter le texte extrait (tronqué si très long)
                if extracted_text:
                    extracted_text = extracted_text.strip()
                    if len(extracted_text) > 12000:
                        extracted_text = extracted_text[:12000] + "\n...[tronqué]"
                    content_parts.append({"type": "input_text", "text": f"Contenu extrait de {file_name}:\n{extracted_text}"})
                elif uploaded_file:
                    # Aucun extrait possible -> fournir un contexte minimal
                    content_parts.append({
                        "type": "input_text",
                        "text": f"Fichier fourni: {file_name}. Si le contenu n'est pas lisible ici, réalise une analyse générique selon l'instruction."
                    })

                system_prompt = (
                    "Tu es un analyste qui extrait des insights d'un document d'entreprise (annuel, durabilité, etc.). "
                    "Tu dois: 1) résumer, 2) extraire des chiffres clés, 3) proposer des visualisations utiles. "
                    "Retourne du JSON valide: { 'summary': str, 'charts': [ { 'title': str, 'type': 'bar|line|pie', 'series': [{ 'name': str, 'x': [...], 'y': [...] }] } ] }"
                )

                user_prompt = "Prépare des graphiques pertinents à partir du rapport ou de l'instruction. Utilise des séries simples et lisibles."

                # Appel IA avec contenu (instruction + document/texte/image)
                resp = client.responses.create(
                    model="gpt-5-mini",
                    input=[
                        {"role": "system", "content": [{"type": "input_text", "text": system_prompt}]},
                        {"role": "user", "content": content_parts + [{"type": "input_text", "text": user_prompt}]}
                    ],
                )

                output = resp.output_text

                # Parsing JSON sécurisé
                import json, re
                txt = output.strip()
                if "```json" in txt:
                    txt = txt.split("```json")[1].split("```", 1)[0]
                elif "```" in txt:
                    txt = txt.split("```", 1)[1]
                txt = re.sub(r"(\d)_(\d)", r"\1\2", txt)
                data = json.loads(txt)

                # Rendu
                st.markdown("### 🧾 Résumé")
                st.write(data.get("summary", "(pas de résumé)"))

                charts = data.get("charts", [])
                if charts:
                    st.markdown("### 📊 Graphiques proposés")
                    for chart in charts:
                        title = chart.get("title", "Graphique")
                        ctype = chart.get("type", "bar")
                        series = chart.get("series", [])
                        fig = go.Figure()
                        for s in series:
                            name = s.get("name", "")
                            x = s.get("x", [])
                            y = s.get("y", [])
                            if ctype == "line":
                                fig.add_trace(go.Scatter(x=x, y=y, mode='lines+markers', name=name))
                            elif ctype == "pie":
                                fig = go.Figure(go.Pie(labels=x, values=y, hole=0.3))
                            else:
                                fig.add_trace(go.Bar(x=x, y=y, name=name))
                        fig.update_layout(title=title, template="plotly_white", height=420)
                        st.plotly_chart(fig, use_container_width=True)
                else:
                    st.info("Aucun graphique proposé par l'IA.")

            except Exception as e:
                st.error(f"Erreur d'analyse: {str(e)}")

def display_company_results(company_data, company_name):
    """Afficher les résultats de l'analyse avec graphiques complets"""
    
    st.markdown(f"## 📊 Résultats pour {company_name}")
    
    # Informations générales
    general_info = company_data.get('general_info', {})
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>🏢 Entreprise</h3>
            <h2>{company_data.get('company_name', 'N/A')}</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>🏭 Secteur</h3>
            <h2>{general_info.get('industry', 'N/A')}</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>📅 Fondée</h3>
            <h2>{general_info.get('founded_year', 'N/A')}</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>🏢 Siège</h3>
            <h2>{general_info.get('headquarters', 'N/A')}</h2>
        </div>
        """, unsafe_allow_html=True)
    
    # Métriques clés 2025
    st.markdown("## 📊 Métriques Clés 2025")
    
    revenue_data = company_data.get('revenue_data', {})
    profit_data = company_data.get('profit_data', {})
    employees_data = company_data.get('employees_data', {})
    esg_data = company_data.get('esg_scores', {})
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        latest_revenue = revenue_data.get('2025', revenue_data.get('2024', 'N/A'))
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>💰 Revenus 2025</h3>
            <h2>${latest_revenue}B</h2>
        </div>
        """ if latest_revenue != 'N/A' else """
        <div class="metric-card-company">
            <h3>💰 Revenus 2025</h3>
            <h2>N/A</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        latest_profit = profit_data.get('2025', profit_data.get('2024', 'N/A'))
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>📈 Bénéfices 2025</h3>
            <h2>${latest_profit}B</h2>
        </div>
        """ if latest_profit != 'N/A' else """
        <div class="metric-card-company">
            <h3>📈 Bénéfices 2025</h3>
            <h2>N/A</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        latest_employees = employees_data.get('2025', employees_data.get('2024', 'N/A'))
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>👥 Employés 2025</h3>
            <h2>{latest_employees:,}</h2>
        </div>
        """ if latest_employees != 'N/A' else """
        <div class="metric-card-company">
            <h3>👥 Employés 2025</h3>
            <h2>N/A</h2>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        latest_esg = esg_data.get('2025', esg_data.get('2024', 'N/A'))
        st.markdown(f"""
        <div class="metric-card-company">
            <h3>📊 Score ESG 2025</h3>
            <h2>{latest_esg}/100</h2>
        </div>
        """ if latest_esg != 'N/A' else """
        <div class="metric-card-company">
            <h3>📊 Score ESG 2025</h3>
            <h2>N/A</h2>
        </div>
        """, unsafe_allow_html=True)
    

    
    # GRAPHIQUES COMPLETS
    st.markdown("## 📈 Analyses Graphiques Complètes")
    
    # Graphiques financiers
    col1, col2 = st.columns(2)
    
    with col1:
        if revenue_data:
            fig_revenue = create_revenue_chart(revenue_data, company_name)
            st.plotly_chart(fig_revenue, use_container_width=True)
    
    with col2:
        if profit_data:
            fig_profit = create_profit_chart(profit_data, company_name)
            st.plotly_chart(fig_profit, use_container_width=True)
    
    # Graphique des employés
    col1, col2 = st.columns(2)
    
    with col1:
        if employees_data:
            fig_employees = create_employees_chart(employees_data, company_name)
            st.plotly_chart(fig_employees, use_container_width=True)
    
    with col2:
        # Espace libre pour équilibrer la mise en page
        st.empty()
    
    # Section Durabilité
    st.markdown("## 🌱 Durabilité & Actualités")
    
    # Actualités de durabilité
    recent_news = company_data.get('recent_news', [])
    sustainability_updates = company_data.get('sustainability_updates', {})
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 📰 Actualités Récentes")
        if recent_news:
            for news in recent_news[:5]:  # Limiter à 5 actualités
                st.markdown(f"• {news}")
        else:
            st.info("Aucune actualité récente disponible")
    
    with col2:
        st.markdown("### 🔄 Mises à jour Durabilité")
        if sustainability_updates:
            if sustainability_updates.get('new_initiatives'):
                st.markdown(f"**🎯 Nouvelles initiatives:** {sustainability_updates['new_initiatives']}")
            if sustainability_updates.get('carbon_progress'):
                st.markdown(f"**🌍 Progrès carbone:** {sustainability_updates['carbon_progress']}")
            if sustainability_updates.get('renewable_projects'):
                st.markdown(f"**⚡ Projets renouvelables:** {sustainability_updates['renewable_projects']}")
            
            # Nouveaux éléments ajoutés
            if sustainability_updates.get('sustainability_target_2030'):
                st.markdown(f"**🎯 Objectifs 2030:** {sustainability_updates['sustainability_target_2030']}")
            if sustainability_updates.get('achievement_sustainability_target_2025'):
                st.markdown(f"**✅ Réalisations 2025:** {sustainability_updates['achievement_sustainability_target_2025']}")
            
            # Liens vers les rapports
            st.markdown("**📄 Rapports officiels:**")
            col_link1, col_link2 = st.columns(2)
            
            with col_link1:
                annual_report_url = sustainability_updates.get('latest_annual_report_url')
                if annual_report_url and annual_report_url != "https://example.com/annual-report-2024":
                    st.markdown(f"[📊 Rapport Annuel]({annual_report_url})")
                else:
                    st.markdown("📊 Rapport Annuel (non disponible)")
            
            with col_link2:
                sustainability_report_url = sustainability_updates.get('latest_sustainability_report_url')
                if sustainability_report_url and sustainability_report_url != "https://example.com/sustainability-report-2024":
                    st.markdown(f"[🌱 Rapport Durabilité]({sustainability_report_url})")
                else:
                    st.markdown("🌱 Rapport Durabilité (non disponible)")
        else:
            st.info("Aucune mise à jour de durabilité disponible")
    
    # Métriques de croissance générale
    st.markdown("### 📈 Croissance & Performance")
    col1, col2 = st.columns(2)
    
    with col1:
        # Graphique de croissance du chiffre d'affaires (taux de croissance)
        if revenue_data and len(revenue_data) > 1:
            fig_growth = create_growth_chart(revenue_data, company_name)
            st.plotly_chart(fig_growth, use_container_width=True)
    
    with col2:
        # Graphique de productivité (revenus par employé)
        if revenue_data and employees_data:
            fig_productivity = create_productivity_chart(revenue_data, employees_data, company_name)
            st.plotly_chart(fig_productivity, use_container_width=True)
    


def create_revenue_chart(revenue_data, company_name):
    """Créer le graphique des revenus"""
    fig = go.Figure()
    
    years = sorted(revenue_data.keys())
    values = [revenue_data[year] for year in years]
    
    fig.add_trace(go.Scatter(
        x=years,
        y=values,
        mode='lines+markers',
        name='Revenus',
        line=dict(color=BLUESIGN_COLORS['primary'], width=4),
        marker=dict(size=10, color=BLUESIGN_COLORS['primary']),
        fill='tonexty'
    ))
    
    fig.update_layout(
        title=f"💰 Évolution des Revenus - {company_name} (2015-2025)",
        xaxis_title="Année",
        yaxis_title="Revenus (milliards $)",
        template="plotly_white",
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def create_profit_chart(profit_data, company_name):
    """Créer le graphique des bénéfices"""
    fig = go.Figure()
    
    years = sorted(profit_data.keys())
    values = [profit_data[year] for year in years]
    
    fig.add_trace(go.Scatter(
        x=years,
        y=values,
        mode='lines+markers',
        name='Bénéfices',
        line=dict(color=BLUESIGN_COLORS['success'], width=4),
        marker=dict(size=10, color=BLUESIGN_COLORS['success']),
        fill='tozeroy'
    ))
    
    fig.update_layout(
        title=f"📈 Évolution des Bénéfices - {company_name} (2015-2025)",
        xaxis_title="Année",
        yaxis_title="Bénéfices (milliards $)",
        template="plotly_white",
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def create_employees_chart(employees_data, company_name):
    """Créer le graphique des employés"""
    fig = go.Figure()
    
    years = sorted(employees_data.keys())
    values = [employees_data[year] for year in years]
    
    fig.add_trace(go.Bar(
        x=years,
        y=values,
        name='Employés',
        marker_color=BLUESIGN_COLORS['accent'],
        marker_line=dict(color='rgba(0,0,0,0.1)', width=1)
    ))
    
    fig.update_layout(
        title=f"👥 Évolution des Employés - {company_name} (2015-2025)",
        xaxis_title="Année",
        yaxis_title="Nombre d'employés",
        template="plotly_white",
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def create_growth_chart(revenue_data, company_name):
    """Créer le graphique de taux de croissance"""
    fig = go.Figure()
    
    years = sorted(revenue_data.keys())
    values = [revenue_data[year] for year in years]
    
    # Calculer les taux de croissance année sur année
    growth_rates = []
    growth_years = []
    
    for i in range(1, len(years)):
        prev_value = values[i-1]
        curr_value = values[i]
        if prev_value > 0:
            growth_rate = ((curr_value - prev_value) / prev_value) * 100
            growth_rates.append(growth_rate)
            growth_years.append(years[i])
    
    # Couleurs conditionnelles (vert pour positif, rouge pour négatif)
    colors = ['green' if rate >= 0 else 'red' for rate in growth_rates]
    
    fig.add_trace(go.Bar(
        x=growth_years,
        y=growth_rates,
        name='Taux de croissance',
        marker_color=colors,
        marker_line=dict(color='rgba(0,0,0,0.3)', width=1)
    ))
    
    fig.update_layout(
        title=f"📈 Taux de Croissance Annuel - {company_name}",
        xaxis_title="Année",
        yaxis_title="Taux de croissance (%)",
        template="plotly_white",
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def create_productivity_chart(revenue_data, employees_data, company_name):
    """Créer le graphique de productivité (revenus par employé)"""
    fig = go.Figure()
    
    # Trouver les années communes
    common_years = sorted(set(revenue_data.keys()) & set(employees_data.keys()))
    
    productivity = []
    for year in common_years:
        if employees_data[year] > 0:
            # Revenus en milliards / employés = revenus par employé en milliers
            prod = (revenue_data[year] * 1000000) / employees_data[year]  # Convertir en dollars par employé
            productivity.append(prod)
        else:
            productivity.append(0)
    
    fig.add_trace(go.Scatter(
        x=common_years,
        y=productivity,
        mode='lines+markers',
        name='Productivité',
        line=dict(color=BLUESIGN_COLORS['success'], width=4),
        marker=dict(size=10, color=BLUESIGN_COLORS['success']),
        fill='tonexty'
    ))
    
    fig.update_layout(
        title=f"💼 Productivité (Revenus/Employé) - {company_name}",
        xaxis_title="Année",
        yaxis_title="Revenus par employé ($)",
        template="plotly_white",
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

if __name__ == "__main__":
    main() 
